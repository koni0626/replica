from __future__ import annotations
"""
services/extract_service.py

アップロード/添付ファイルからテキストを抽出する簡易サービス。
フェーズ2: Office/PDF対応（docx, pptx, xlsx, pdf）。

依存ライブラリ（requirements.txtに追記済み）
- python-docx
- python-pptx
- openpyxl
- pdfminer.six

注意:
- 抽出はテキスト化のみ。バイナリのプレビューは対象外。
- PDFはレイアウトが崩れることがある。OCRは別途実装（オプション）。
"""

from pathlib import Path
from typing import Optional


def _safe_str(s: Optional[str]) -> str:
    return s or ""


class ExtractService:
    @staticmethod
    def extract_text(path: str | Path, ext: Optional[str] = None, limit: Optional[int] = None) -> str:
        """
        与えられたファイルからテキストを抽出して返す。
        - ext が指定されない場合は拡張子から自動判定
        - limit が指定されれば、その文字数（近似）で切り詰め
        """
        p = Path(path)
        if not p.exists() or not p.is_file():
            return ""
        ext = (ext or p.suffix.lstrip(".")).lower()

        text = ""
        try:
            if ext in {"txt", "md", "markdown", "csv", "json", "yaml", "yml", "html", "htm",
                       "py", "js", "ts", "java", "php", "go", "rb", "cs", "sh", "sql", "css"}:
                text = p.read_text(encoding="utf-8", errors="ignore")
            elif ext == "docx":
                text = ExtractService._extract_docx(p)
            elif ext == "pptx":
                text = ExtractService._extract_pptx(p)
            elif ext == "xlsx":
                text = ExtractService._extract_xlsx(p)
            elif ext == "pdf":
                text = ExtractService._extract_pdf(p)
            else:
                # 未対応拡張子は空文字
                text = ""
        except Exception:
            text = ""

        if limit is not None and limit > 0:
            return text[:limit]
        return text

    @staticmethod
    def _extract_docx(p: Path) -> str:
        try:
            from docx import Document
        except Exception:
            return ""
        try:
            doc = Document(str(p))
            parts = []
            for para in doc.paragraphs:
                parts.append(para.text)
            return "\n".join(filter(None, parts))
        except Exception:
            return ""

    @staticmethod
    def _extract_pptx(p: Path) -> str:
        try:
            from pptx import Presentation
        except Exception:
            return ""
        try:
            prs = Presentation(str(p))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text"):
                            parts.append(_safe_str(shape.text))
                    except Exception:
                        continue
            return "\n".join(filter(None, parts))
        except Exception:
            return ""

    @staticmethod
    def _extract_xlsx(p: Path) -> str:
        try:
            import openpyxl
        except Exception:
            return ""
        try:
            wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append(f"# Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    # タブ区切りで1行に整形
                    cells = ["" if v is None else str(v) for v in row]
                    parts.append("\t".join(cells))
            return "\n".join(parts)
        except Exception:
            return ""

    @staticmethod
    def _extract_pdf(p: Path) -> str:
        try:
            from pdfminer.high_level import extract_text
        except Exception:
            return ""
        try:
            return extract_text(str(p))
        except Exception:
            return ""
