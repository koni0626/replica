"""
PowerPoint (.pptx) 読み取り用ツール

- 依存: python-pptx
    pip install python-pptx

- 役割:
  * 指定パスの .pptx ファイルからスライド内のテキストを抽出して返す
  * タイトル/本文/表のセル/グループ化された図形内のテキストを可能な範囲で取得
  * スピーカーノート（講演者ノート）もオプションで抽出

- 公開関数:
    read_pptx_text(path, limit_chars=None, include_slide_header=True,
                   include_notes=True, table_cell_delimiter="\t") -> str

- 使用例:
    from tools.office_pptx_tool import read_pptx_text
    text = read_pptx_text(r"C:\\path\\to\\file.pptx", limit_chars=5000)

注意:
- .ppt (97-2003) 形式は対象外です。必要に応じて .pptx へ変換してください。
- 画像内の文字（OCR）は未対応です。
"""
from __future__ import annotations

import os
from typing import Iterable, Optional

try:
    from pptx import Presentation  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
except Exception:
    Presentation = None  # type: ignore
    MSO_SHAPE_TYPE = None  # type: ignore


class PowerPointReadError(Exception):
    """PowerPoint 読み取り時の一般例外"""


def _is_pptx(path: str) -> bool:
    return os.path.splitext(path)[1].lower() == ".pptx"


def _iter_shape_text(shape) -> Iterable[str]:
    """shape からテキストを（可能な限り）抽出して列挙する。

    - 通常のテキストフレーム（shape.has_text_frame）
    - 表（shape.has_table）: セルをタブ区切りで連結して1行として返す
    - グループ図形（GROUP）: 再帰的に処理
    """
    # テキストフレーム
    if getattr(shape, "has_text_frame", False):
        try:
            txt = shape.text
            if txt:
                yield txt
        except Exception:
            pass

    # 表
    if getattr(shape, "has_table", False):
        try:
            tbl = shape.table
            for row in tbl.rows:
                cells = []
                for cell in row.cells:
                    try:
                        cells.append(cell.text or "")
                    except Exception:
                        cells.append("")
                yield "\t".join(cells)
        except Exception:
            pass

    # グループ図形
    try:
        if MSO_SHAPE_TYPE is not None and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                yield from _iter_shape_text(sub)
    except Exception:
        pass


def read_pptx_text(
    path: str,
    limit_chars: Optional[int] = None,
    include_slide_header: bool = True,
    include_notes: bool = True,
    table_cell_delimiter: str = "\t",
) -> str:
    """
    .pptx からテキストを抽出して返す。

    Args:
        path: 対象ファイルの絶対 or 相対パス
        limit_chars: 返却する最大文字数（None なら上限なし）。超過時は末尾を切り詰め。
        include_slide_header: True の場合、各スライドの先頭に "# Slide: <n> <title>" を付与
        include_notes: True の場合、スピーカーノートを "## Notes" セクションとして付与
        table_cell_delimiter: 表のセルを連結する区切り文字（既定: タブ）

    Returns:
        抽出テキスト（失敗時は例外）

    Raises:
        PowerPointReadError: 依存未導入 / ファイル不正 / 読み取り失敗など致命的なときに発生
    """
    if not _is_pptx(path):
        raise PowerPointReadError(f"Unsupported extension for PPTX reader: {path}")

    if Presentation is None:
        # 依存がない環境では例外
        raise PowerPointReadError(
            "python-pptx is not installed. Please `pip install python-pptx`."
        )

    if not os.path.exists(path):
        raise PowerPointReadError(f"File not found: {path}")

    try:
        prs = Presentation(path)
        parts: list[str] = []

        for idx, slide in enumerate(prs.slides, start=1):
            # スライドヘッダ
            if include_slide_header:
                title = None
                try:
                    if slide.shapes.title:
                        title = slide.shapes.title.text
                except Exception:
                    title = None
                header = f"# Slide {idx}: {title}" if title else f"# Slide {idx}"
                parts.append(header)

            # スライド内の図形テキスト
            for shp in slide.shapes:
                for line in _iter_shape_text(shp):
                    if table_cell_delimiter != "\t":
                        # _iter_shape_text は表のセル連結にタブを使用しているため置換で合わせる
                        line = line.replace("\t", table_cell_delimiter)
                    parts.append(line)

            # スピーカーノート
            if include_notes:
                try:
                    notes_slide = getattr(slide, "notes_slide", None)
                    if notes_slide and notes_slide.notes_text_frame:
                        notes = notes_slide.notes_text_frame.text or ""
                        if notes.strip():
                            parts.append("## Notes")
                            parts.append(notes)
                except Exception:
                    pass

        text = "\n".join(parts)
        if limit_chars is not None and limit_chars > 0 and len(text) > limit_chars:
            return text[:limit_chars]
        return text

    except PowerPointReadError:
        raise
    except Exception as e:
        raise PowerPointReadError(f"Failed to read pptx: {e}") from e
